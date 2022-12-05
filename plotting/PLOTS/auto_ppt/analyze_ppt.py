#QUICK START GUIDE
#https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#add-picture-example
#https://pbpython.com/creating-powerpoint.html

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
 

title='Abhishek Nath\n\
Dept. of Physics and Astronomy, PI, Uni Heidelberg  '
Alicelogo='ALICE_logo.png'
HDlogo='Heidelberg.png'
Highrrlogo='HighRR_logo.png'
PIlogo='PI_logo.png'
pylogo='ALICE3.jpg'
pptlogo='ALICE3.jpg'
prs = Presentation()

# front page
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)



shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(9/1.5),Inches(16),Inches(9/8.5)
)
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= title
line=shape.line
line.color.rgb=RGBColor(68,114,148)
#logo1=slide.shapes.add_picture(pylogo,Inches(13.8),Inches(6.0),height=Inches(1.0),width=Inches(1.0))
logo2=slide.shapes.add_picture(Alicelogo,Inches(14.0),Inches(6.3),height=Inches(1.5),width=Inches(1.5))
logo3=slide.shapes.add_picture(HDlogo,Inches(12.0),Inches(6.3),height=Inches(1.5),width=Inches(1.5))
logo4=slide.shapes.add_picture(Highrrlogo,Inches(10.0),Inches(6.3),height=Inches(1.5),width=Inches(1.5))
logo4=slide.shapes.add_picture(PIlogo,Inches(8.0),Inches(6.3),height=Inches(1.5),width=Inches(1.5))
#-----------------------------------------------------------------------------------------------------------------------

#loc1='/home/abhishek/PhD/Work/work_A/photons/anaConv/plots/pp_0.1pTcut/hDiffRap_ratio_Pi0.png'
#loc2='/home/abhishek/PhD/Work/work_A/photons/anaConv/plots/pp_0.2pTcut/hDiffRap_ratio_Pi0.png'
#loc3='/home/abhishek/PhD/Work/work_A/photons/anaConv/plots/pp_0.3pTcut/hDiffRap_ratio_Pi0.png'
#loc4='/home/abhishek/PhD/Work/work_A/photons/anaConv/plots/pp_0.5pTcut/hDiffRap_ratio_Pi0.png'
#loc5='/home/abhishek/PhD/Work/work_A/photons/anaConv/plots/pp_0.1pTcut/hDiffRap_ratio_Pi0.png'

#loc1='./hInvMassGGB.png'

loc1a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_MassResolution/PbPb/meson_MassResolution_Pi0_Barrel.png'
loc1b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_MassResolution/PbPb/meson_MassResolution_Pi0_Forward.png'
loc2a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_MassResolution/PbPb/meson_MassResolution_Eta_Barrel.png'
loc2b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_MassResolution/PbPb/meson_MassResolution_Eta_Forward.png'
loc3a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_Acceptance/PbPb/Acceptance_meson_ALL.png'
loc3b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_Effeciency/PbPb/meson_Effeciency.png'
loc4a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_NetEffeciency/PbPb/meson_NetEffeciency.png'
loc4b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_dNdyPythia/PbPb/meson_dNdyPythia.png'
loc5a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_openingangle/PbPb/meson_openingangle_Pi0.png'
loc5b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_openingangle/PbPb/meson_openingangle_Eta.png'
loc6='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Pi0_B/MesonSubtractedPureGaussianFit.png'
loc7='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Pi0_B/meson_signalandbackground.png'
loc8='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Pi0_F/MesonSubtractedPureGaussianFit.png'
loc9='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Pi0_F/meson_signalandbackground.png'
loc10='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Eta_B/MesonSubtractedPureGaussianFit.png'
loc11='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Eta_B/meson_signalandbackground.png'
loc12='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Eta_F/MesonSubtractedPureGaussianFit.png'
loc13='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Eta_F/meson_signalandbackground.png'

loc14a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Pi0_F/meson_FIT_mu.png'
loc14b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Eta_F/meson_FIT_mu.png'
loc14c='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Pi0_B/meson_FIT_mu.png'
loc14d='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Eta_B/meson_FIT_mu.png'


loc15a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Pi0_F/meson_FIT_sigma.png'
loc15b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Eta_F/meson_FIT_sigma.png'
loc15c='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Pi0_B/meson_FIT_sigma.png'
loc15d='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_ptbins_fit/PbPb/Eta_B/meson_FIT_sigma.png'

loc16a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_signal_AND_background/PbPb/meson_signal.png'
loc16b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_signal_AND_background/PbPb/meson_background.png'

loc17a='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_signalVSbackground/PbPb/meson_signalVSbackground.png'
loc17b='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS/meson_significance/PbPb//meson_significance_ALL.png'



#Page 1
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Mass resolution Pi0 "

pic1a = slide.shapes.add_picture(loc1a, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pic1b = slide.shapes.add_picture(loc1b, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



#Page 2
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Mass resolution Eta "
pic2a = slide.shapes.add_picture(loc2a, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pic2b = slide.shapes.add_picture(loc2b, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))


#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



 #Page 3
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Acceptance and Efficiency "

pic3a = slide.shapes.add_picture(loc3a, Inches(.5), Inches(1.75),width=Inches(7), height=Inches(7))
pic3b = slide.shapes.add_picture(loc3b, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



 #Page 4
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Net Efficiency and PYTHIA dN/dy"

pic4a = slide.shapes.add_picture(loc4a, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pic4b = slide.shapes.add_picture(loc4b, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



 #Page 5
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Opening angle"

pic5a = slide.shapes.add_picture(loc5a, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pic5b = slide.shapes.add_picture(loc5b, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



 #Page 6
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal Fit Pi0 B"
pic6 = slide.shapes.add_picture(loc6, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


 #Page 7
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal and Background comparison Pi0 B"
pic7 = slide.shapes.add_picture(loc7, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True




 #Page 8
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal Fit Pi0 F"
pic8 = slide.shapes.add_picture(loc8, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





 #Page 9
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal and Background comparison Pi0 F"
pic9 = slide.shapes.add_picture(loc9, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





 #Page 10
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal Fit Eta B"
pic10 = slide.shapes.add_picture(loc10, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



 #Page 11
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal and Background comparison Eta B"
pic11 = slide.shapes.add_picture(loc11, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True

 #Page 12
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal Fit Eta F"
pic12 = slide.shapes.add_picture(loc12, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



 #Page 13
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal and Background comparison Eta F"
pic13 = slide.shapes.add_picture(loc13, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



 #Page 14
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Gaussian fitted mass"

pic14a = slide.shapes.add_picture(loc14a, Inches(1.5), Inches(1.75),width=Inches(6), height=Inches(3))
pic14b = slide.shapes.add_picture(loc14b, Inches(1.5), Inches(4.75),width=Inches(6), height=Inches(3))

pic14c = slide.shapes.add_picture(loc14c, Inches(8.5), Inches(1.75),width=Inches(6), height=Inches(3))
pic14d = slide.shapes.add_picture(loc14d, Inches(8.5), Inches(4.75),width=Inches(6), height=Inches(3))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



 #Page 15
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Gaussian fitted sigma"

pic15a = slide.shapes.add_picture(loc15a, Inches(1.5), Inches(1.75),width=Inches(6), height=Inches(3))
pic15b = slide.shapes.add_picture(loc15b, Inches(1.5), Inches(4.75),width=Inches(6), height=Inches(3))

pic15c = slide.shapes.add_picture(loc15c, Inches(8.5), Inches(1.75),width=Inches(6), height=Inches(3))
pic15d = slide.shapes.add_picture(loc15d, Inches(8.5), Inches(4.75),width=Inches(6), height=Inches(3))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





 #Page 16
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Net Background and Net Signal"

pic16a = slide.shapes.add_picture(loc16a, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pic16b = slide.shapes.add_picture(loc16b, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True






 #Page 17
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: signal to background ratio and Significance"

pic17a = slide.shapes.add_picture(loc17a, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pic17b = slide.shapes.add_picture(loc17b, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





prs.save('meson_results.pptx')
