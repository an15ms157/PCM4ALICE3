#QUICK START GUIDE
#https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#add-picture-example
#https://pbpython.com/creating-powerpoint.html
# HELPER: https://www.shibutan-bloomers.com/python-libraly-pptx_en/6132/

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
from datetime import date


Author = 'Abhishek Nath'
Dept = 'Dept. of Physics and Astronomy'
Uni ='PI, Uni Heidelberg'

today = date.today()

Alicelogo='ALICE_logo.png'
HDlogo='Heidelberg.png'
Highrrlogo='HighRR_logo.png'
PIlogo='PI_logo.png'
pylogo='ALICE3.jpg'
pptlogo='ALICE3.jpg'
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

BASE_LOC='./'
PbPb_MassResolution_Pi0_Barrel_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_MassResolution/PbPb/meson_MassResolution_Pi0_Barrel.png'
PbPb_MassResolution_Pi0_Forward_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_MassResolution/PbPb/meson_MassResolution_Pi0_Forward.png'
PbPb_MassResolution_Eta_Barrel_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_MassResolution/PbPb/meson_MassResolution_Eta_Barrel.png'
PbPb_MassResolution_Eta_Forward_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_MassResolution/PbPb/meson_MassResolution_Eta_Forward.png'

pp_MassResolution_Pi0_Barrel_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_MassResolution/pp/meson_MassResolution_Pi0_Barrel.png'
pp_MassResolution_Pi0_Forward_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_MassResolution/pp/meson_MassResolution_Pi0_Forward.png'
pp_MassResolution_Eta_Barrel_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_MassResolution/pp/meson_MassResolution_Eta_Barrel.png'
pp_MassResolution_Eta_Forward_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_MassResolution/pp/meson_MassResolution_Eta_Forward.png'



PbPb_Acceptance_meson_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_Acceptance/PbPb/Acceptance_meson_ALL.png'
PbPb_meson_Effeciency_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_Effeciency/PbPb/meson_Effeciency.png'
pp_Acceptance_meson_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_Acceptance/pp/Acceptance_meson_ALL.png'
pp_meson_Effeciency_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_Effeciency/pp/meson_Effeciency.png'

PbPb_NetEffeciency_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_NetEffeciency/PbPb/meson_NetEffeciency.png'
PbPb_dNdyPythia_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_dNdyPythia/PbPb/meson_dNdyPythia.png'
pp_NetEffeciency_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_NetEffeciency/pp/meson_NetEffeciency.png'
pp_dNdyPythia_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_dNdyPythia/pp/meson_dNdyPythia.png'

PbPb_openingangle_Pi0_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_openingangle/PbPb/meson_openingangle_Pi0.png'
PbPb_openingangle_Eta_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_openingangle/PbPb/meson_openingangle_Eta.png'
pp_openingangle_Pi0_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_openingangle/pp/meson_openingangle_Pi0.png'
pp_openingangle_Eta_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_openingangle/pp/meson_openingangle_Eta.png'

PbPb_meson_signalandbackground_Pi0B_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Pi0_B/meson_signalandbackground.png'
PbPb_meson_signalandbackground_EtaB_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Eta_B/meson_signalandbackground.png'
PbPb_meson_signalandbackground_Pi0F_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Pi0_F/meson_signalandbackground.png'
PbPb_meson_signalandbackground_EtaF_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Eta_F/meson_signalandbackground.png'

PbPb_MesonSubtractedPureGaussianFit_Pi0B_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Pi0_B/MesonSubtractedPureGaussianFit.png'
PbPb_MesonSubtractedPureGaussianFit_EtaB_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Eta_B/MesonSubtractedPureGaussianFit.png'
PbPb_MesonSubtractedPureGaussianFit_Pi0F_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Pi0_F/MesonSubtractedPureGaussianFit.png'
PbPb_MesonSubtractedPureGaussianFit_EtaF_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Eta_F/MesonSubtractedPureGaussianFit.png'

pp_meson_signalandbackground_Pi0B_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Pi0_B/meson_signalandbackground.png'
pp_meson_signalandbackground_EtaB_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Eta_B/meson_signalandbackground.png'
pp_meson_signalandbackground_Pi0F_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Pi0_F/meson_signalandbackground.png'
pp_meson_signalandbackground_EtaF_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Eta_F/meson_signalandbackground.png'

pp_MesonSubtractedPureGaussianFit_Pi0B_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Pi0_B/MesonSubtractedPureGaussianFit.png'
pp_MesonSubtractedPureGaussianFit_EtaB_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Eta_B/MesonSubtractedPureGaussianFit.png'
pp_MesonSubtractedPureGaussianFit_Pi0F_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Pi0_F/MesonSubtractedPureGaussianFit.png'
pp_MesonSubtractedPureGaussianFit_EtaF_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Eta_F/MesonSubtractedPureGaussianFit.png'

PbPb_FIT_mu_Pi0B_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Pi0_B/meson_FIT_mu.png'
PbPb_FIT_mu_Pi0F_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Pi0_F/meson_FIT_mu.png'
PbPb_FIT_mu_EtaB_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Eta_B/meson_FIT_mu.png'
PbPb_FIT_mu_EtaF_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Eta_F/meson_FIT_mu.png'

pp_FIT_mu_Pi0B_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Pi0_B/meson_FIT_mu.png'
pp_FIT_mu_Pi0F_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Pi0_F/meson_FIT_mu.png'
pp_FIT_mu_EtaB_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Eta_B/meson_FIT_mu.png'
pp_FIT_mu_EtaF_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Eta_F/meson_FIT_mu.png'

PbPb_FIT_sigma_Pi0B_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Pi0_B/meson_FIT_sigma.png'
PbPb_FIT_sigma_Pi0F_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Pi0_F/meson_FIT_sigma.png'
PbPb_FIT_sigma_EtaB_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Eta_B/meson_FIT_sigma.png'
PbPb_FIT_sigma_EtaF_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/PbPb/Eta_F/meson_FIT_sigma.png'

pp_FIT_sigma_Pi0B_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Pi0_B/meson_FIT_sigma.png'
pp_FIT_sigma_Pi0F_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Pi0_F/meson_FIT_sigma.png'
pp_FIT_sigma_EtaB_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Eta_B/meson_FIT_sigma.png'
pp_FIT_sigma_EtaF_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_ptbins_fit/pp/Eta_F/meson_FIT_sigma.png'

PbPb_background_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_signal_AND_background/PbPb/meson_background.png'
PbPb_signal_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_signal_AND_background/PbPb/meson_signal.png'

pp_background_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_signal_AND_background/pp/meson_background.png'
pp_signal_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_signal_AND_background/pp/meson_signal.png'

PbPb_signalVSbackground_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_signal_VS_background/PbPb/meson_signalVSbackground.png'
PbPb_significance_ALL_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_significance/PbPb/meson_significance_ALL.png'

pp_signalVSbackground_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_signal_VS_background/pp/meson_signalVSbackground.png'
pp_significance_ALL_loc='/home/abhishek/PhD/Work/work_A/photons/PCM4ALICE3/plotting/PLOTS1/meson_significance/pp/meson_significance_ALL.png'



comversion_vertices_B_loc="/home/abhishek/PhD/Work/work_A/photons/auto_ppt/figures/vertex_F.png"
comversion_vertices_F_loc="/home/abhishek/PhD/Work/work_A/photons/auto_ppt/figures/Conversion_vertices.png"

Cocktail_contribution_loc="./figures/CocktailGammasRatioToAll_0_80_10210a13_0d200009ab770c00amd0404000_0152101500000000.png"
PCM_loc="/home/abhishek/PhD/Work/work_A/photons/auto_ppt/figures/PCM.png"
ALICE3_loc="/home/abhishek/PhD/Work/work_A/photons/auto_ppt/figures/ALICE3.png"



# front page
#-----------------------------------------------------------------------------------------------------------------------
layout6 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout6)

#TITLE OPF THE PRESENTATION
txBox00 = slide.shapes.add_textbox(Inches(3), Inches(1), width=Inches(7), height=Inches(2))
tf = txBox00.text_frame
p01 = tf.add_paragraph()
p01.text = "ALICE 3 light meson update" 
p01.font.name = "Serif"
p01.font.size = Pt(48)
p01.font.bold = True

txBox01 = slide.shapes.add_textbox(Inches(1), Inches(3), width=Inches(7), height=Inches(7))
tf = txBox01.text_frame

p02 = tf.add_paragraph()
p02.font.name = "Serif"
p02.font.size = Pt(30)
p02.text = Author

p03 = tf.add_paragraph()
p03.font.name = "Serif"
p03.font.size = Pt(30)
p03.text = Dept 

p04 = tf.add_paragraph()
p04.font.name = "Serif"
p04.font.size = Pt(30)
p04.text = Uni 


txBox02 = slide.shapes.add_textbox(Inches(5), Inches(7.5), width=Inches(7), height=Inches(1))
tf = txBox02.text_frame

p05 = tf.add_paragraph()
p05.font.name = "Serif"
p05.font.size = Pt(24)
p05.text = "18th Nov 2022 : Update on Systematic studies"

#logo1=slide.shapes.add_picture(pylogo,Inches(13.8),Inches(6.0),height=Inches(1.0),width=Inches(1.0))
logo2=slide.shapes.add_picture(Alicelogo,Inches(14.0),Inches(0.5),height=Inches(1.5),width=Inches(1.5))
logo3=slide.shapes.add_picture(HDlogo,Inches(14.0),Inches(2.5),height=Inches(1.5),width=Inches(1.5))
logo4=slide.shapes.add_picture(Highrrlogo,Inches(14.0),Inches(4.0),height=Inches(1.5),width=Inches(1.5))
logo4=slide.shapes.add_picture(PIlogo,Inches(14.0),Inches(5.5),height=Inches(1.5),width=Inches(1.5))





#Page 1
#-----------------------------------------------------------------------------------------------------------------------
#layout3 = prs.slide_layouts[3]  #TWO COLUMNS
slide = prs.slides.add_slide(layout6) 
shapes = slide.shapes

###############   Title:

left = Inches(1.0)
top = Inches(0.5) 
width = Inches(2.0)
height = Inches(1.0)
txBox10 = slide.shapes.add_textbox(left, 0, width=Inches(7), height=Inches(6.5))
tf = txBox10.text_frame
p10 = tf.add_paragraph()
p10.font.name = "Serif"
p10.font.size = Pt(32)
p10.font.bold = True
p10.text = "Neutral meson Monte-Carlo \nperformance studies using PCM"


##################  Table:

rows = 4
cols = 3
table = shapes.add_table(rows, cols, left, 9*top, width, height).table

# set column widths
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.5)

# write column headings
table.cell(0, 1).text = 'Barrel'
table.cell(0, 2).text = 'Forward'

# write row headings
table.cell(1, 0).text = 'Rapidity'
table.cell(2, 0).text = 'Conversion'
table.cell(3, 0).text = 'Reconstruction'

# write row headings
table.cell(1, 1).text = '0<|y|<1.3 '
table.cell(1, 2).text = '1.75<|y|<4 '

table.cell(2, 1).text = 'R<22 cm '
table.cell(2, 2).text = '0<Z<135 cm '

table.cell(3, 1).text = '5 layers \n22<R<100 cm '
table.cell(3, 2).text = '5 disks \n135<|Z|<200 cm '



################## BULLETS:

txBox11 = slide.shapes.add_textbox(left, 13*top, 3*width, 3*height)
tf = txBox11.text_frame

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- Full simulations of pp √s = 14 and \n PbPb √s = 5.52 TeV using PYTHIA 8.2\n"

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- pp: 9.6∗10^6 ; PbPb: 19k events"


left_arrow = width + left + Inches(0.15)
top_arrow = 5*top+ height/2 
height_arrow = Inches(0.15)
width_arrow = Inches(0.75)
del_x = Inches(0.15)

################## FLOW CHARTS:

box1=slide.shapes.add_shape(    MSO_SHAPE.ROUNDED_RECTANGLE, left, 5*top, width, height)
pg1 = box1.text_frame.paragraphs[0] 
pg1.text = 'Event Generator'         	
pg1.font.size = Pt(18)    

slide.shapes.add_shape(    MSO_SHAPE.RIGHT_ARROW, left_arrow, top_arrow, width_arrow, height_arrow)
txBox12 = slide.shapes.add_textbox(left+del_x, top_arrow, width=Inches(4.5), height=Inches(4))
tf = txBox12.text_frame
p12 = tf.add_paragraph()
p12.font.name = "Serif"
p12.font.size = Pt(18)
p12.text = "\tFull \nsimulations"


box2 = slide.shapes.add_shape(    MSO_SHAPE.ROUNDED_RECTANGLE, left+Inches(3.0), 5*top, width, height)
pg2 = box2.text_frame.paragraphs[0] 
pg2.text = 'Photon Performance'         	
pg2.font.size = Pt(18) 

slide.shapes.add_shape(    MSO_SHAPE.RIGHT_ARROW, left_arrow +Inches(3.0) , top_arrow, width_arrow, height_arrow)
txBox13 = slide.shapes.add_textbox(left+width_arrow+ width +3*del_x, top_arrow, width=Inches(4.5), height=Inches(4))
tf = txBox13.text_frame
p13 = tf.add_paragraph()
p13.font.name = "Serif"
p13.font.size = Pt(18)
p13.text = "\tFast \nsimulations"

box3 = slide.shapes.add_shape(    MSO_SHAPE.ROUNDED_RECTANGLE, left+Inches(6.0), 5*top, width, height)
pg3 = box3.text_frame.paragraphs[0] 
pg3.text = 'Meson Performance'         	
pg3.font.size = Pt(18)    

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


################   FIGURE: 
comversion_vertices_B_fig = slide.shapes.add_picture(comversion_vertices_B_loc, 10*left, 2*top, 2*width, 4*height)
comversion_vertices_F_fig = slide.shapes.add_picture(comversion_vertices_B_loc, 10*left, 10*top, 2*width, 4*height)







#Page 2
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 
shapes = slide.shapes

###############   Title:

left = Inches(1.0)
top = Inches(0.5) 
width = Inches(2.0)
height = Inches(1.0)
txBox10 = slide.shapes.add_textbox(left, 0, width=Inches(7), height=Inches(6.5))
tf = txBox10.text_frame
p10 = tf.add_paragraph()
p10.font.name = "Serif"
p10.font.size = Pt(32)
p10.font.bold = True
p10.text = "ALICE 3 and recent developments"

################## BULLETS:
txBox11 = slide.shapes.add_textbox(8*left, 3*top, 2.5*width, 4*height)
tf = txBox11.text_frame


p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- ALICE 3: An ultra-light silicon wafer\n based detector system with very\n low momentum tracking capabilities for\n high luminosity collisions\n"

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- High-resolution vertex detector,\n PID over large acceptance\n"

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- Photons (1 MeV/c - 50 GeV/c), \nwide rapidity range (|y | < 4) :\n → ECAL\n → FCT\n → PCM\n"


left_arrow = width + left + Inches(0.15)
top_arrow = 5*top+ height/2 
height_arrow = Inches(0.15)
width_arrow = Inches(0.75)
del_x = Inches(0.15)

################   FIGURE: 
ALICE3_fig = slide.shapes.add_picture(ALICE3_loc,.5*left, 3*top, 2.75*width, 4*height)





#Page 3
#-----------------------------------------------------------------------------------------------------------------------

slide = prs.slides.add_slide(layout6) 
shapes = slide.shapes

###############   Title:

left = Inches(1.0)
top = Inches(0.5) 
width = Inches(2.0)
height = Inches(1.0)
txBox10 = slide.shapes.add_textbox(left, 0, width=Inches(7), height=Inches(6.5))
tf = txBox10.text_frame
p10 = tf.add_paragraph()
p10.font.name = "Serif"
p10.font.size = Pt(32)
p10.font.bold = True
p10.text = "PCM & Direct Photons"

################## BULLETS:
txBox11 = slide.shapes.add_textbox(9*left, 12*top, 2.5*width, 4*height)
tf = txBox11.text_frame

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- Can escape QCD medium without \n being affected\n"

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- Carry information on medium’s space-time expansion \nand temperature: hadron matter phase diagram\n"

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- High precision decay photon measurement necessary. \nMain contributions from π0 and η mesons"


txBox11 = slide.shapes.add_textbox(.5*left, 2*top, 2.5*width, 4*height)
tf = txBox11.text_frame

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- PCM :Converted photons measured by \nreconstructing e+e− pairs using a V0\n algorithm.\n"

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- Photon momentum resolution \nlinked to charged particle \nmomentum resolution.\n"

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- Good measurement at low p_{T}\n"

p11 = tf.add_paragraph()
p11.font.name = "Serif"
p11.font.size = Pt(18)
p11.text = "- π0 and η are reconstructed \nin their respective 2γ decay channel\n"

left_arrow = width + left + Inches(0.15)
top_arrow = 5*top+ height/2 
height_arrow = Inches(0.15)
width_arrow = Inches(0.75)
del_x = Inches(0.15)

################   FIGURE: 
PCM_loc_fig = slide.shapes.add_picture(PCM_loc, 2*left, 10*top , 2.75*width, 3.5*height)
Cocktail_contribution_fig = slide.shapes.add_picture(Cocktail_contribution_loc, 9*left, 1.5*top, 2.25*width, 5*height)


                    #############   PbPb    ###############

slide = prs.slides.add_slide(layout6) 

txBox10 = slide.shapes.add_textbox(left+Inches(6), Inches(4), width=Inches(7), height=Inches(6.5))
tf = txBox10.text_frame
p10 = tf.add_paragraph()
p10.font.name = "Serif"
p10.font.size = Pt(32)
p10.font.bold = True
p10.text = "PbPb results"







#MASS RESOLUTION 
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Mass resolution Pi0 "
PbPb_MassResolution_Pi0_Barrel_loc = slide.shapes.add_picture(PbPb_MassResolution_Pi0_Barrel_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
PbPb_MassResolution_Pi0_Forward_loc = slide.shapes.add_picture(PbPb_MassResolution_Pi0_Forward_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Mass resolution Eta "
PbPb_MassResolution_Eta_Barrel_loc = slide.shapes.add_picture(PbPb_MassResolution_Eta_Barrel_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
PbPb_MassResolution_Eta_Forward_loc = slide.shapes.add_picture(PbPb_MassResolution_Eta_Forward_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True






# ACCEPTANCE AND EFFICIENCY
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Acceptance and Efficiency "

PbPb_Acceptance_meson_fig = slide.shapes.add_picture(PbPb_Acceptance_meson_loc, Inches(.5), Inches(1.75),width=Inches(7), height=Inches(7))
PbPb_meson_Effeciency_fig = slide.shapes.add_picture(PbPb_meson_Effeciency_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



#NET EFFICIENCY AND pYTHIA dNDy
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Net Efficiency and PYTHIA dN/dy"

PbPb_NetEffeciency_fig = slide.shapes.add_picture(PbPb_NetEffeciency_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
PbPb_dNdyPythia_fig = slide.shapes.add_picture(PbPb_dNdyPythia_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True









#OPENING ANGLE
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Opening angle"

PbPb_openingangle_Pi0_fig = slide.shapes.add_picture(PbPb_openingangle_Pi0_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
PbPb_openingangle_Eta_fig = slide.shapes.add_picture(PbPb_openingangle_Eta_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True







#GAUSSIAN PEAK FIT 
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal Fit Pi0 B"
PbPb_MesonSubtractedPureGaussianFit_Pi0B_fig = slide.shapes.add_picture(PbPb_MesonSubtractedPureGaussianFit_Pi0B_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal Fit Pi0 F"
PbPb_MesonSubtractedPureGaussianFit_Pi0F_fig = slide.shapes.add_picture(PbPb_MesonSubtractedPureGaussianFit_Pi0F_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal Fit Eta B"
PbPb_MesonSubtractedPureGaussianFit_EtaB_fig = slide.shapes.add_picture(PbPb_MesonSubtractedPureGaussianFit_EtaB_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Signal Fit Eta F"
PbPb_MesonSubtractedPureGaussianFit_EtaF_fig = slide.shapes.add_picture(PbPb_MesonSubtractedPureGaussianFit_EtaF_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True






#SIGNAL AND BACKGROUND COMPARISON
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results Pi0 B: Signal and Background comparison"
PbPb_meson_signalandbackground_Pi0B_fig = slide.shapes.add_picture(PbPb_meson_signalandbackground_Pi0B_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True

#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results Pi0 F: Signal and Background comparison"
PbPb_meson_signalandbackground_Pi0F_fig = slide.shapes.add_picture(PbPb_meson_signalandbackground_Pi0F_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True



#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results Eta B: Signal and Background comparison"
PbPb_meson_signalandbackground_EtaB_fig = slide.shapes.add_picture(PbPb_meson_signalandbackground_EtaB_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results Eta F: Signal and Background comparison"
PbPb_meson_signalandbackground_EtaF_fig = slide.shapes.add_picture(PbPb_meson_signalandbackground_EtaF_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





#GAUSSIAN MEAN FIT
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Gaussian fitted mass"

PbPb_FIT_mu_Pi0B_fig = slide.shapes.add_picture(PbPb_FIT_mu_Pi0B_loc, Inches(1.5), Inches(1.75),width=Inches(6), height=Inches(3))
PbPb_FIT_mu_Pi0F_fig = slide.shapes.add_picture(PbPb_FIT_mu_Pi0F_loc, Inches(1.5), Inches(4.75),width=Inches(6), height=Inches(3))

PbPb_FIT_mu_EtaB_fig = slide.shapes.add_picture(PbPb_FIT_mu_EtaB_loc, Inches(8.5), Inches(1.75),width=Inches(6), height=Inches(3))
PbPb_FIT_mu_EtaF_fig = slide.shapes.add_picture(PbPb_FIT_mu_EtaF_loc, Inches(8.5), Inches(4.75),width=Inches(6), height=Inches(3))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





#GAUSSIAN SIGMA FIT
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Gaussian fitted sigma"

PbPb_FIT_sigma_Pi0B_fig = slide.shapes.add_picture(PbPb_FIT_sigma_Pi0B_loc, Inches(1.5), Inches(1.75),width=Inches(6), height=Inches(3))
PbPb_FIT_sigma_Pi0F_fig = slide.shapes.add_picture(PbPb_FIT_sigma_Pi0F_loc, Inches(1.5), Inches(4.75),width=Inches(6), height=Inches(3))

PbPb_FIT_sigma_EtaB_fig = slide.shapes.add_picture(PbPb_FIT_sigma_EtaB_loc, Inches(8.5), Inches(1.75),width=Inches(6), height=Inches(3))
PbPb_FIT_sigma_EtaF_fig = slide.shapes.add_picture(PbPb_FIT_sigma_EtaF_loc, Inches(8.5), Inches(4.75),width=Inches(6), height=Inches(3))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True







# NET SIGNAL AND BACKGROUND
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: Net Background and Net Signal"

PbPb_background_loc = slide.shapes.add_picture(PbPb_background_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
PbPb_signal_loc = slide.shapes.add_picture(PbPb_signal_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





#SIGNAL TO BACKGROUND AND SIGNIFICANCE
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "PbPb Results: signal to background ratio and Significance"

PbPb_signalVSbackground_fig = slide.shapes.add_picture(PbPb_signalVSbackground_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
PbPb_significance_ALL_fig = slide.shapes.add_picture(PbPb_significance_ALL_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True










                    #############   pp    ###############

slide = prs.slides.add_slide(layout6) 

txBox10 = slide.shapes.add_textbox(left+Inches(6), Inches(4), width=Inches(7), height=Inches(6.5))
tf = txBox10.text_frame
p10 = tf.add_paragraph()
p10.font.name = "Serif"
p10.font.size = Pt(32)
p10.font.bold = True
p10.text = "pp results"



#MASS RESOLUTION 
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Mass resolution Pi0 "
pp_MassResolution_Pi0_Barrel_loc = slide.shapes.add_picture(pp_MassResolution_Pi0_Barrel_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pp_MassResolution_Pi0_Forward_loc = slide.shapes.add_picture(pp_MassResolution_Pi0_Forward_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Mass resolution Eta "
pp_MassResolution_Eta_Barrel_loc = slide.shapes.add_picture(pp_MassResolution_Eta_Barrel_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pp_MassResolution_Eta_Forward_loc = slide.shapes.add_picture(pp_MassResolution_Eta_Forward_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True








# ACCEPTANCE AND EFFICIENCY
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Acceptance and Efficiency "

pp_Acceptance_meson_fig = slide.shapes.add_picture(pp_Acceptance_meson_loc, Inches(.5), Inches(1.75),width=Inches(7), height=Inches(7))
pp_meson_Effeciency_fig = slide.shapes.add_picture(pp_meson_Effeciency_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True






#NET EFFICIENCY AND pYTHIA dNDy
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Net Efficiency and PYTHIA dN/dy"

pp_NetEffeciency_fig = slide.shapes.add_picture(pp_NetEffeciency_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pp_dNdyPythia_fig = slide.shapes.add_picture(pp_dNdyPythia_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True






#OPENING ANGLE
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 


#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Opening angle"

pp_openingangle_Pi0_fig = slide.shapes.add_picture(pp_openingangle_Pi0_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pp_openingangle_Eta_fig = slide.shapes.add_picture(pp_openingangle_Eta_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True






#GAUSSIAN PEAK FIT 
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Signal Fit Pi0 B"
pp_MesonSubtractedPureGaussianFit_Pi0B_fig = slide.shapes.add_picture(pp_MesonSubtractedPureGaussianFit_Pi0B_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True

#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Signal Fit Pi0 F"
pp_MesonSubtractedPureGaussianFit_Pi0F_fig = slide.shapes.add_picture(pp_MesonSubtractedPureGaussianFit_Pi0F_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Signal Fit Eta B"
pp_MesonSubtractedPureGaussianFit_EtaB_fig = slide.shapes.add_picture(pp_MesonSubtractedPureGaussianFit_EtaB_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Signal Fit Eta F"
pp_MesonSubtractedPureGaussianFit_EtaF_fig = slide.shapes.add_picture(pp_MesonSubtractedPureGaussianFit_EtaF_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True







#SIGNAL AND BACKGROUND COMPARISON
#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results Pi0 B: Signal and Background comparison"
pp_meson_signalandbackground_Pi0B_fig = slide.shapes.add_picture(pp_meson_signalandbackground_Pi0B_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results Pi0 F: Signal and Background comparison"
pp_meson_signalandbackground_Pi0F_fig = slide.shapes.add_picture(pp_meson_signalandbackground_Pi0F_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results Eta B: Signal and Background comparison"
pp_meson_signalandbackground_EtaB_fig = slide.shapes.add_picture(pp_meson_signalandbackground_EtaB_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True


#-----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(layout6) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results Eta F: Signal and Background comparison"
pp_meson_signalandbackground_EtaF_fig = slide.shapes.add_picture(pp_meson_signalandbackground_EtaF_loc, Inches(1.0), Inches(1.75),width=Inches(14), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True






#GAUSSIAN MEAN FIT
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Gaussian fitted mass"

pp_FIT_mu_Pi0B_fig = slide.shapes.add_picture(pp_FIT_mu_Pi0B_loc, Inches(1.5), Inches(1.75),width=Inches(6), height=Inches(3))
pp_FIT_mu_Pi0F_fig = slide.shapes.add_picture(pp_FIT_mu_Pi0F_loc, Inches(1.5), Inches(4.75),width=Inches(6), height=Inches(3))

pp_FIT_mu_EtaB_fig = slide.shapes.add_picture(pp_FIT_mu_EtaB_loc, Inches(8.5), Inches(1.75),width=Inches(6), height=Inches(3))
pp_FIT_mu_EtaF_fig = slide.shapes.add_picture(pp_FIT_mu_EtaF_loc, Inches(8.5), Inches(4.75),width=Inches(6), height=Inches(3))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





#GAUSSIAN SIGMA FIT
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Gaussian fitted sigma"

pp_FIT_sigma_Pi0B_fig = slide.shapes.add_picture(pp_FIT_sigma_Pi0B_loc, Inches(1.5), Inches(1.75),width=Inches(6), height=Inches(3))
pp_FIT_sigma_Pi0F_fig = slide.shapes.add_picture(pp_FIT_sigma_Pi0F_loc, Inches(1.5), Inches(4.75),width=Inches(6), height=Inches(3))

pp_FIT_sigma_EtaB_fig = slide.shapes.add_picture(pp_FIT_sigma_EtaB_loc, Inches(8.5), Inches(1.75),width=Inches(6), height=Inches(3))
pp_FIT_sigma_EtaF_fig = slide.shapes.add_picture(pp_FIT_sigma_EtaF_loc, Inches(8.5), Inches(4.75),width=Inches(6), height=Inches(3))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True






# NET SIGNAL AND BACKGROUND
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: Net Background and Net Signal"

pp_background_loc = slide.shapes.add_picture(pp_background_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pp_signal_loc = slide.shapes.add_picture(pp_signal_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





#SIGNAL TO BACKGROUND AND SIGNIFICANCE
#-----------------------------------------------------------------------------------------------------------------------
layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 

#IMAGE
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.5))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(68,114,148)
shape.text= "pp Results: signal to background ratio and Significance"

pp_signalVSbackground_fig = slide.shapes.add_picture(pp_signalVSbackground_loc, Inches(0.5), Inches(1.75),width=Inches(7), height=Inches(7))
pp_significance_ALL_fig = slide.shapes.add_picture(pp_significance_ALL_loc, Inches(8.0), Inches(1.75),width=Inches(7), height=Inches(7))

#PAGE NUMBER
slide_id = prs.slides.index(slide)
txBox = slide.shapes.add_textbox(Inches(15), Inches(8), width=Inches(1), height=Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "%d" % slide_id 
p.font.bold = True





 #Page 17
#-----------------------------------------------------------------------------------------------------------------------


# TEST PAGE 
#-----------------------------------------------------------------------------------------------------------------------

layout8 = prs.slide_layouts[6]
slide = prs.slides.add_slide(layout8) 


#PAGE NUMBER
#for slide in prs.slides:
#    slide_id = prs.slides.index(slide)
#    txBox = slide.shapes.add_textbox(Inches(14), Inches(8), width=Inches(1), height=Inches(1))
#    tf = txBox.text_frame
#    p = tf.add_paragraph()
#    p.text = "%d" % slide_id 
#    p.font.bold = True

prs.save('meson_results.pptx')
